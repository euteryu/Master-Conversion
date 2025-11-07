# Save this file as backend/backend_ai_generator.py
import ollama
from pptx import Presentation
from pptx.util import Inches
import os
import re

# Define the model to be used by Ollama. 'llama3:8b' is recommended.
# LLM_MODEL = 'llama3:8b'
# LLM_MODEL = 'qwen3:8b'
LLM_MODEL = 'phi3:mini'

def _create_llm_prompt(user_prompt, output_format):
    """
    Wraps the user's prompt with specific instructions for the LLM 
    to ensure structured, predictable output.
    """
    if output_format == 'powerpoint':
        return f"""
        You are an expert presentation creator. Based on the following topic, generate content for a 5-slide presentation.
        The first slide should be a title slide with a suitable title and a brief subtitle.
        The remaining slides should have a title and content with 3-5 concise bullet points.
        Your entire response MUST strictly follow this format, with no extra commentary:

        SLIDE 1 (Title Slide)
        TITLE: [Your Title Here]
        SUBTITLE: [Your Subtitle Here]

        SLIDE 2
        TITLE: [Your Title Here]
        CONTENT:
        - Bullet point 1.
        - Bullet point 2.
        - Bullet point 3.

        SLIDE 3
        TITLE: [Your Title Here]
        CONTENT:
        - Bullet point 1.
        - Bullet point 2.
        - Bullet point 3.
        
        ... and so on.

        TOPIC: "{user_prompt}"
        """
    # Add logic for 'poster' and 'flashcard' formats here in the future
    else:
        # Fallback for now
        return user_prompt

def _parse_llm_response(response_text):
    """
    Parses the raw text from the LLM into a structured list of slide data.
    """
    slides = []
    # Split the response into individual slide sections
    slide_chunks = response_text.strip().split('SLIDE ')
    
    for chunk in slide_chunks:
        if not chunk.strip():
            continue
            
        slide_data = {}
        
        # Use regex to find TITLE, SUBTITLE, and CONTENT
        title_match = re.search(r'TITLE:\s*(.*)', chunk, re.IGNORECASE)
        subtitle_match = re.search(r'SUBTITLE:\s*(.*)', chunk, re.IGNORECASE)
        content_match = re.search(r'CONTENT:\s*([\s\S]*)', chunk, re.IGNORECASE)

        if title_match:
            slide_data['title'] = title_match.group(1).strip()
        
        if subtitle_match:
            slide_data['subtitle'] = subtitle_match.group(1).strip()
            
        if content_match:
            # Clean up content, remove leading dashes and extra whitespace
            content = content_match.group(1).strip()
            lines = [line.strip().lstrip('- ').strip() for line in content.split('\n')]
            slide_data['content'] = "\n".join(lines)
            
        if slide_data:
            slides.append(slide_data)
            
    return slides

def _create_ppt_from_content(content_list, theme, output_path, templates_dir):
    """
    Generates a .pptx file from the structured content using a specified template.
    This version is more robust and finds placeholders by type, not by a fixed index.
    """
    theme_to_template = {
        'Revision Theme': 'revision_theme.pptx',
        'Minutes Theme': 'minutes_theme.pptx',
        'Upload Custom Theme': 'revision_theme.pptx' # Default for now
    }
    
    template_name = theme_to_template.get(theme, 'revision_theme.pptx')
    template_path = os.path.join(templates_dir, template_name)

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")

    prs = Presentation(template_path)

    # Assuming layout indices: 0 for Title Slide, 1 for Title and Content
    # These indices are standard in most PowerPoint themes.
    title_slide_layout = prs.slide_layouts[0] 
    content_slide_layout = prs.slide_layouts[1]

    for i, slide_content in enumerate(content_list):
        if i == 0 and 'subtitle' in slide_content: # Title Slide
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            # Find the subtitle placeholder by its type to be safe
            subtitle = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 'SUBTITLE' or shape.placeholder_format.type == 'BODY':
                    subtitle = shape
                    break
            if title:
                title.text = slide_content.get('title', 'Title')
            if subtitle:
                subtitle.text = slide_content.get('subtitle', '')
        else: # Content Slide
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            
            # --- THIS IS THE CRITICAL FIX ---
            # Instead of assuming placeholders[1], find the main body placeholder
            body_shape = None
            for shape in slide.placeholders:
                # A content placeholder is usually of type 'BODY'
                if shape.placeholder_format.type == 'BODY':
                    body_shape = shape
                    break
            
            if title:
                title.text = slide_content.get('title', 'Content Title')
            
            # Only try to add text if we actually found the body placeholder
            if body_shape:
                tf = body_shape.text_frame
                tf.clear() # Good practice to clear existing placeholder text
                p = tf.paragraphs[0]
                p.text = slide_content.get('content', '')
            else:
                print(f"Warning: Could not find a 'BODY' placeholder on the content slide layout for theme '{theme}'.")
            
    prs.save(output_path)


def generate_presentation(prompt, output_format, theme, output_path, templates_dir, progress_callback):
    """
    The main orchestrator function for generating a presentation.
    """
    try:
        # Step 1: Create the detailed prompt for the LLM
        progress_callback(1, 4, 'Engineering AI prompt...')
        engineered_prompt = _create_llm_prompt(prompt, output_format)

        # Step 2: Communicate with the local LLM
        progress_callback(2, 4, 'Generating content with local AI...')
        response = ollama.chat(
            model=LLM_MODEL,
            messages=[{'role': 'user', 'content': engineered_prompt}]
        )
        raw_content = response['message']['content']

        # Step 3: Parse the LLM's response into a usable structure
        progress_callback(3, 4, 'Parsing AI response...')
        structured_content = _parse_llm_response(raw_content)

        if not structured_content:
            return "AI failed to generate content in the expected format. Please try again."

        # Step 4: Build the PowerPoint file
        progress_callback(4, 4, 'Building PowerPoint file...')
        if output_format == 'powerpoint':
            _create_ppt_from_content(structured_content, theme, output_path, templates_dir)
        else:
            return f"Output format '{output_format}' is not yet supported."
            
        return None # Success
        
    except FileNotFoundError as e:
        return f"Error: A required template file was not found. Details: {e}"
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        import traceback
        traceback.print_exc()
        return f"An unexpected error occurred during generation: {e}"