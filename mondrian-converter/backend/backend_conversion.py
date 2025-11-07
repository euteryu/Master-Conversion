# Save this file as backend_conversion.py (FIXED VERSION)

import fitz
from pptx import Presentation
from pptx.util import Pt, Inches
import io
import traceback

POINTS_TO_EMUS = 12700

def convert_pdf_to_ppt_hybrid(pdf_path, ppt_path, progress_callback, dpi=150):
    """
    Converts PDF to a hybrid PPT using a robust in-memory redaction method.
    The 'dpi' parameter controls the quality and speed of background rendering.
    """
    try:
        original_doc = fitz.open(pdf_path)
        if original_doc.page_count == 0:
            return "The selected PDF is empty."
        
        # Extract text data from all pages
        original_text_data = []
        for page in original_doc:
            try:
                text_dict = page.get_text("dict")
                original_text_data.append(text_dict)
            except Exception as e:
                print(f"Warning: Could not extract text from page {page.number}: {e}")
                original_text_data.append({"blocks": []})
        
        # Create clean document by redacting text
        clean_doc_in_memory = fitz.open()
        for page_num, page in enumerate(original_doc):
            try:
                # Get text data for this page
                if page_num < len(original_text_data):
                    page_text_data = original_text_data[page_num]
                else:
                    page_text_data = {"blocks": []}
                
                # Redact text blocks
                for block in page_text_data.get("blocks", []):
                    if block.get('type') == 0:  # Text block
                        for line in block.get('lines', []):
                            for span in line.get('spans', []):
                                if 'bbox' in span:
                                    redaction_rect = fitz.Rect(span['bbox'])
                                    redaction_rect.x0 -= 2
                                    redaction_rect.y0 -= 2
                                    redaction_rect.x1 += 2
                                    redaction_rect.y1 += 2
                                    page.add_redact_annot(redaction_rect)
                
                # Apply redactions
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                
                # Insert page WITHOUT copying links/annotations to avoid index errors
                clean_doc_in_memory.insert_pdf(
                    original_doc, 
                    from_page=page.number, 
                    to_page=page.number,
                    links=False,  # Don't copy links
                    annots=False  # Don't copy annotations
                )
            except Exception as e:
                print(f"Warning: Could not process page {page_num}: {e}")
                # Try inserting without links/annotations
                try:
                    clean_doc_in_memory.insert_pdf(
                        original_doc, 
                        from_page=page.number, 
                        to_page=page.number,
                        links=False,
                        annots=False
                    )
                except Exception as e2:
                    print(f"Error: Could not insert page {page_num} at all: {e2}")
                    # Create a blank page as fallback
                    clean_doc_in_memory.new_page(
                        width=page.rect.width,
                        height=page.rect.height
                    )
        
        original_doc.close()

        # Create PowerPoint presentation
        prs = Presentation()
        
        # Set slide dimensions based on first page
        if len(clean_doc_in_memory) > 0:
            first_page = clean_doc_in_memory.load_page(0)
            prs.slide_width = int(first_page.rect.width * POINTS_TO_EMUS)
            prs.slide_height = int(first_page.rect.height * POINTS_TO_EMUS)

        # Process each page
        for page_num, page in enumerate(clean_doc_in_memory):
            try:
                # Render page to image
                pix = page.get_pixmap(dpi=dpi)
                
                # Create slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(
                    io.BytesIO(pix.tobytes("png")), 0, 0,
                    width=prs.slide_width, height=prs.slide_height
                )

                # Add text overlays if available
                if page_num < len(original_text_data):
                    page_text_data = original_text_data[page_num]
                    
                    for block in page_text_data.get("blocks", []):
                        if block.get('type') == 0:  # Text block
                            for line in block.get('lines', []):
                                for span in line.get('spans', []):
                                    try:
                                        # Validate span has required fields
                                        if not all(key in span for key in ['bbox', 'text', 'size']):
                                            continue
                                        
                                        x0, y0, x1, y1 = span['bbox']
                                        
                                        # Add text box
                                        txBox = slide.shapes.add_textbox(
                                            int(x0 * POINTS_TO_EMUS), 
                                            int(y0 * POINTS_TO_EMUS),
                                            int((x1 - x0) * POINTS_TO_EMUS), 
                                            int((y1 - y0) * POINTS_TO_EMUS)
                                        )
                                        
                                        txBox.fill.background()
                                        txBox.line.fill.background()
                                        
                                        p = txBox.text_frame.paragraphs[0]
                                        p.margin_left = 0
                                        p.margin_right = 0
                                        p.margin_top = 0
                                        p.margin_bottom = 0
                                        
                                        run = p.add_run()
                                        run.text = span.get('text', '')
                                        
                                        font = run.font
                                        font.size = Pt(int(span.get('size', 12)))
                                        
                                        # Check for bold font
                                        if 'font' in span and "bold" in span['font'].lower():
                                            font.bold = True
                                    except Exception as e:
                                        print(f"Warning: Could not add text span on page {page_num}: {e}")
                                        continue
                
                progress_callback(page_num + 1, len(clean_doc_in_memory))
            except Exception as e:
                print(f"Warning: Error processing page {page_num}: {e}")
                progress_callback(page_num + 1, len(clean_doc_in_memory))
                continue

        prs.save(ppt_path)
        clean_doc_in_memory.close()
        return None
    except Exception as e:
        traceback.print_exc()
        return f"Hybrid conversion failed: {str(e)}"


def convert_ppt_to_pdf(ppt_path, pdf_path, progress_callback):
    """
    Converts a PowerPoint file (.ppt or .pptx) to PDF.
    Requires Microsoft PowerPoint to be installed on the system.
    """
    powerpoint = None
    presentation = None
    try:
        progress_callback(1, 2, "Initializing PowerPoint...")
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1  # Keep it visible for debugging, change to 0 for hidden

        # Ensure absolute paths are used, as COM objects can be sensitive
        ppt_path_abs = os.path.abspath(ppt_path)
        pdf_path_abs = os.path.abspath(pdf_path)

        if not os.path.exists(ppt_path_abs):
            return f"Input file not found: {ppt_path_abs}"

        presentation = powerpoint.Presentations.Open(ppt_path_abs)
        
        # ppSaveAsPDF is format type 32
        progress_callback(2, 2, "Saving as PDF...")
        presentation.SaveAs(pdf_path_abs, 32)
        
        return None # Success
    except Exception as e:
        traceback.print_exc()
        return f"PPT to PDF conversion failed. Ensure PowerPoint is installed. Error: {e}"
    finally:
        if presentation:
            presentation.Close()
        if powerpoint:
            powerpoint.Quit()


CONVERSION_STRATEGIES = {
    'pdf_to_ppt': convert_pdf_to_ppt_hybrid,
    'ppt_to_pdf': convert_ppt_to_pdf,
}