# In backend_conversion.py

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt, Inches
import io
import traceback

POINTS_TO_EMUS = 12700  # PDF points to PPTX EMUs

def convert_pdf_to_ppt_hybrid(pdf_path, ppt_path, progress_callback):
    """
    Converts PDF to a hybrid PPT using an expanded redaction method to remove
    tight-fitting background shapes behind text.
    """
    try:
        # Step 1: Get original text content and positions
        original_doc = fitz.open(pdf_path)
        if original_doc.page_count == 0:
            return "The selected PDF is empty."
        
        original_text_data = [page.get_text("dict") for page in original_doc]
        
        # Step 2: Create a text-free version of the PDF in memory
        clean_doc_in_memory = fitz.open()
        for page_num, page in enumerate(original_doc):
            # Mark every text span for redaction
            for block in original_text_data[page_num]["blocks"]:
                if block['type'] == 0:
                    for line in block['lines']:
                        for span in line['spans']:
                            # --- THE FIX: Expand the redaction area ---
                            # Create a rectangle from the text's bounding box
                            redaction_rect = fitz.Rect(span['bbox'])
                            # Expand it by 2 pixels in each direction to catch background boxes
                            redaction_rect.x0 -= 2
                            redaction_rect.y0 -= 2
                            redaction_rect.x1 += 2
                            redaction_rect.y1 += 2
                            page.add_redact_annot(redaction_rect)
            
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE) # Only remove text/vectors, not images
            clean_doc_in_memory.insert_pdf(original_doc, from_page=page.number, to_page=page.number)
        
        original_doc.close()

        # Step 3: Build the PowerPoint
        prs = Presentation()
        first_page = clean_doc_in_memory.load_page(0)
        prs.slide_width = int(first_page.rect.width * POINTS_TO_EMUS)
        prs.slide_height = int(first_page.rect.height * POINTS_TO_EMUS)

        for page_num, page in enumerate(clean_doc_in_memory):
            pix = page.get_pixmap(dpi=150)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")), 0, 0,
                width=prs.slide_width, height=prs.slide_height
            )

            page_text_data = original_text_data[page_num]
            for block in page_text_data["blocks"]:
                if block['type'] == 0:
                    for line in block['lines']:
                        for span in line['spans']:
                            x0, y0, x1, y1 = span['bbox']
                            txBox = slide.shapes.add_textbox(
                                int(x0 * POINTS_TO_EMUS), int(y0 * POINTS_TO_EMUS),
                                int((x1 - x0) * POINTS_TO_EMUS), int((y1 - y0) * POINTS_TO_EMUS)
                            )
                            txBox.fill.background()
                            txBox.line.fill.background()
                            
                            p = txBox.text_frame.paragraphs[0]
                            p.margin_left = 0
                            p.margin_right = 0
                            p.margin_top = 0
                            p.margin_bottom = 0

                            run = p.add_run()
                            run.text = span['text']
                            font = run.font
                            font.size = Pt(int(span['size']))
                            if "bold" in span['font'].lower(): font.bold = True
            
            progress_callback(page_num + 1, len(clean_doc_in_memory))

        prs.save(ppt_path)
        clean_doc_in_memory.close()
        return None
    except Exception as e:
        traceback.print_exc()
        return f"Hybrid conversion failed: {str(e)}"

# --- The rest of the file (convert_pdf_to_ppt_image_only and CONVERSION_STRATEGIES) remains unchanged ---
def convert_pdf_to_ppt_image_only(pdf_path, ppt_path, progress_callback):
    try:
        pdf_doc = fitz.open(pdf_path)
        prs = Presentation()
        total_pages = len(pdf_doc)
        if total_pages == 0: return "The selected PDF is empty."
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        for page_num, page in enumerate(pdf_doc):
            pix = page.get_pixmap(dpi=150)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")), Inches(0), Inches(0),
                width=prs.slide_width, height=prs.slide_height
            )
            progress_callback(page_num + 1, total_pages)
        prs.save(ppt_path)
        pdf_doc.close()
        return None
    except Exception as e:
        traceback.print_exc()
        return f"Image-only conversion failed: {str(e)}"

CONVERSION_STRATEGIES = {
    ('PDF', 'PPT', 'Hybrid (Editable Text)'): convert_pdf_to_ppt_hybrid,
    ('PDF', 'PPT', 'Image Only (Flattened)'): convert_pdf_to_ppt_image_only,
}